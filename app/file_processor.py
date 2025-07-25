import os
import io
from PIL import Image
from pptx import Presentation
import PyPDF2
from docx import Document
import tempfile

# 可选导入 - 如果依赖包不可用，功能会被禁用
try:
    import cv2
    OPENCV_AVAILABLE = True
except ImportError:
    OPENCV_AVAILABLE = False
    cv2 = None

try:
    import easyocr
    EASYOCR_AVAILABLE = True
except ImportError:
    EASYOCR_AVAILABLE = False
    easyocr = None

try:
    import speech_recognition as sr
    SPEECH_RECOGNITION_AVAILABLE = True
except ImportError:
    SPEECH_RECOGNITION_AVAILABLE = False
    sr = None

try:
    from moviepy.editor import VideoFileClip
    MOVIEPY_AVAILABLE = True
except ImportError:
    MOVIEPY_AVAILABLE = False
    VideoFileClip = None

try:
    from pydub import AudioSegment
    PYDUB_AVAILABLE = True
except ImportError:
    PYDUB_AVAILABLE = False
    AudioSegment = None

class FileProcessor:
    def __init__(self):
        self.ocr_reader = None
        if EASYOCR_AVAILABLE:
            try:
                self.ocr_reader = easyocr.Reader(['ch_sim', 'en'])
            except Exception as e:
                print(f"警告：EasyOCR 初始化失败: {e}")
                self.ocr_reader = None
        
    def process_file(self, file_path, content_type):
        """
        根据文件类型处理文件并提取文本内容
        """
        try:
            if content_type == 'text':
                return self.extract_text_from_txt(file_path)
            elif content_type == 'ppt':
                return self.extract_text_from_ppt(file_path)
            elif content_type == 'pdf':
                return self.extract_text_from_pdf(file_path)
            elif content_type == 'audio':
                return self.extract_text_from_audio(file_path)
            elif content_type == 'video':
                return self.extract_text_from_video(file_path)
            else:
                raise ValueError(f"不支持的文件类型: {content_type}")
        except Exception as e:
            print(f"处理文件时出错: {str(e)}")
            return ""
    
    def extract_text_from_txt(self, file_path):
        """从文本文件提取内容"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # 尝试其他编码
            with open(file_path, 'r', encoding='gbk') as file:
                return file.read()
    
    def extract_text_from_ppt(self, file_path):
        """从PowerPoint文件提取文本内容"""
        text_content = []
        presentation = Presentation(file_path)
        
        for slide in presentation.slides:
            # 提取文本框内容
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
            
            # 提取图片中的文字（OCR）
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        pil_image = Image.open(io.BytesIO(image_bytes))
                        
                        # 使用OCR提取图片中的文字（如果可用）
                        if self.ocr_reader is not None:
                            ocr_result = self.ocr_reader.readtext(image_bytes)
                            for detection in ocr_result:
                                text_content.append(detection[1])
                        else:
                            text_content.append("[图片内容 - OCR不可用]")
                    except Exception as e:
                        print(f"OCR处理图片时出错: {str(e)}")
                        continue
        
        return '\n'.join(text_content)
    
    def extract_text_from_pdf(self, file_path):
        """从PDF文件提取文本内容"""
        text_content = []
        
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    if text.strip():
                        text_content.append(text)
                    
                    # 如果文本提取失败，尝试OCR
                    if not text.strip():
                        try:
                            # 这里需要将PDF页面转换为图片然后OCR
                            # 暂时跳过，需要额外的库如pdf2image
                            pass
                        except Exception as e:
                            print(f"OCR处理PDF页面时出错: {str(e)}")
                            continue
        except Exception as e:
            print(f"读取PDF文件时出错: {str(e)}")
        
        return '\n'.join(text_content)
    
    def extract_text_from_audio(self, file_path):
        """从音频文件提取文本内容（语音识别）"""
        if not SPEECH_RECOGNITION_AVAILABLE or not PYDUB_AVAILABLE:
            return "[音频文字 - 语音识别库不可用]"
            
        recognizer = sr.Recognizer()
        
        try:
            # 将音频文件转换为wav格式
            audio = AudioSegment.from_file(file_path)
            
            # 创建临时wav文件
            with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as temp_wav:
                audio.export(temp_wav.name, format="wav")
                
                # 使用语音识别
                with sr.AudioFile(temp_wav.name) as source:
                    audio_data = recognizer.record(source)
                    
                # 尝试识别（中文和英文）
                try:
                    text = recognizer.recognize_google(audio_data, language='zh-CN')
                except sr.UnknownValueError:
                    try:
                        text = recognizer.recognize_google(audio_data, language='en-US')
                    except sr.UnknownValueError:
                        text = ""
                
                # 清理临时文件
                os.unlink(temp_wav.name)
                
                return text
                
        except Exception as e:
            print(f"音频识别时出错: {str(e)}")
            return ""
    
    def extract_text_from_video(self, file_path):
        """从视频文件提取文本内容（音频转文字 + OCR画面文字）"""
        if not MOVIEPY_AVAILABLE:
            return "[视频文字 - MoviePy库不可用]"
            
        text_content = []
        
        try:
            # 加载视频
            video = VideoFileClip(file_path)
            
            # 提取音频并转换为文字
            if video.audio:
                with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as temp_audio:
                    video.audio.write_audiofile(temp_audio.name)
                    audio_text = self.extract_text_from_audio(temp_audio.name)
                    if audio_text:
                        text_content.append(f"音频内容: {audio_text}")
                    os.unlink(temp_audio.name)
            
            # 提取关键帧并OCR识别文字
            duration = video.duration
            num_frames = min(10, int(duration))  # 最多提取10帧
            
            for i in range(num_frames):
                timestamp = (duration / num_frames) * i
                frame = video.get_frame(timestamp)
                
                # 将帧转换为PIL图像
                pil_image = Image.fromarray(frame)
                
                # 将PIL图像转换为字节
                img_byte_arr = io.BytesIO()
                pil_image.save(img_byte_arr, format='PNG')
                img_byte_arr = img_byte_arr.getvalue()
                
                # OCR识别（如果可用）
                try:
                    if self.ocr_reader is not None:
                        ocr_result = self.ocr_reader.readtext(img_byte_arr)
                        frame_text = []
                        for detection in ocr_result:
                            if detection[2] > 0.5:  # 置信度阈值
                                frame_text.append(detection[1])
                    else:
                        frame_text = ["[视频帧文字 - OCR不可用]"]
                    
                    if frame_text:
                        text_content.append(f"画面文字 ({timestamp:.1f}s): {' '.join(frame_text)}")
                except Exception as e:
                    print(f"OCR处理视频帧时出错: {str(e)}")
                    continue
            
            video.close()
            
        except Exception as e:
            print(f"处理视频文件时出错: {str(e)}")
        
        return '\n'.join(text_content)
    
    def extract_text_from_docx(self, file_path):
        """从Word文档提取文本内容"""
        try:
            doc = Document(file_path)
            text_content = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            return '\n'.join(text_content)
        except Exception as e:
            print(f"处理Word文档时出错: {str(e)}")
            return ""
    
    def extract_text_from_pdf_bytes(self, pdf_bytes):
        """从PDF字节流提取文本内容"""
        text_content = []
        
        try:
            pdf_file = io.BytesIO(pdf_bytes)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            print(f"PDF文件包含 {len(pdf_reader.pages)} 页")
            
            for page_num in range(len(pdf_reader.pages)):
                try:
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    if text.strip():
                        text_content.append(text)
                        print(f"第{page_num+1}页提取到 {len(text)} 字符")
                    else:
                        print(f"第{page_num+1}页未能提取到文本")
                except Exception as e:
                    print(f"处理第{page_num+1}页时出错: {e}")
                    continue
                    
        except Exception as e:
            print(f"从PDF字节流读取文本时出错: {str(e)}")
            return f"PDF文件处理失败: {str(e)}"
        
        result = '\n'.join(text_content)
        print(f"PDF处理完成，总共提取 {len(result)} 字符")
        
        if not result.strip():
            return "PDF文件中未能提取到文本内容，可能是扫描版PDF或图片格式"
        
        return result

    def extract_text_from_ppt_bytes(self, ppt_bytes):
        """从PowerPoint字节流提取文本内容"""
        text_content = []
        
        try:
            ppt_file = io.BytesIO(ppt_bytes)
            
            # 检查文件是否为有效的PowerPoint文件
            if len(ppt_bytes) < 100:
                return "文件太小，可能不是有效的PPT文件"
            
            # 检查文件头以确定文件类型
            ppt_file.seek(0)
            header = ppt_file.read(8)
            ppt_file.seek(0)
            
            print(f"PPT文件大小: {len(ppt_bytes)} 字节")
            print(f"文件头: {header.hex()}")
            
            # .pptx 文件是ZIP格式，应该以 'PK' 开头
            # .ppt 文件是OLE格式，有不同的头部
            if header[:2] == b'PK':
                # 这是 .pptx 文件 (ZIP格式)
                try:
                    presentation = Presentation(ppt_file)
                    slide_count = len(presentation.slides)
                    print(f"PPTX文件包含 {slide_count} 张幻灯片")
                    
                    for slide_num, slide in enumerate(presentation.slides):
                        slide_text = []
                        # 提取文本框内容
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text.append(shape.text.strip())
                        
                        if slide_text:
                            slide_content = '\n'.join(slide_text)
                            text_content.append(slide_content)
                            print(f"第{slide_num+1}张幻灯片提取到 {len(slide_content)} 字符")
                        else:
                            print(f"第{slide_num+1}张幻灯片未找到文本内容")
                            
                except Exception as e:
                    print(f"处理PPTX文件时出错: {str(e)}")
                    # 如果python-pptx处理失败，尝试其他方法
                    return self._extract_text_from_ppt_alternative(ppt_bytes)
            
            elif header[:8] == b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1':
                # 这是老版本的 .ppt 文件 (OLE格式)
                print("检测到旧版本PPT文件(.ppt)")
                return "不支持旧版本PPT文件，请将文件转换为.pptx格式后重新上传"
            
            else:
                # 未知文件格式
                print(f"未知的文件格式，文件头: {header.hex()}")
                return "文件格式不正确，请确保上传的是有效的PPT或PPTX文件"
                
        except Exception as e:
            print(f"从PPT字节流读取文本时出错: {str(e)}")
            return f"PPT文件处理失败: {str(e)}"
        
        result = '\n'.join(text_content)
        print(f"PPT处理完成，总共提取 {len(result)} 字符")
        
        if not result.strip():
            return "PPT文件中未能提取到文本内容，请检查文件是否包含文本"
        
        return result
    
    def _extract_text_from_ppt_alternative(self, ppt_bytes):
        """PPT文件处理的备用方法"""
        try:
            # 尝试将文件保存为临时文件然后处理
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
                temp_file.write(ppt_bytes)
                temp_file.flush()
                
                # 使用文件路径方式处理
                text_content = self.extract_text_from_ppt(temp_file.name)
                
                # 清理临时文件
                os.unlink(temp_file.name)
                
                return text_content
        except Exception as e:
            print(f"备用PPT处理方法也失败: {str(e)}")
            return f"PPT文件处理失败: {str(e)}"

# 文件类型检测函数
def detect_file_type(filename):
    """根据文件扩展名检测文件类型"""
    ext = os.path.splitext(filename)[1].lower()
    
    type_mapping = {
        '.txt': 'text',
        '.md': 'text',
        '.ppt': 'ppt',
        '.pptx': 'ppt',
        '.pdf': 'pdf',
        '.doc': 'text',
        '.docx': 'text',
        '.mp3': 'audio',
        '.wav': 'audio',
        '.m4a': 'audio',
        '.flac': 'audio',
        '.mp4': 'video',
        '.avi': 'video',
        '.mov': 'video',
        '.wmv': 'video',
        '.flv': 'video'
    }
    
    return type_mapping.get(ext, 'unknown')
