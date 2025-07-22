#!/usr/bin/env python3
"""
创建一个简单的PPTX文件用于测试
"""

from pptx import Presentation
from pptx.util import Inches

def create_test_pptx():
    """创建测试用的PPTX文件"""
    print("创建测试PPTX文件...")
    
    # 创建演示文稿
    prs = Presentation()
    
    # 添加第一张幻灯片
    slide_layout = prs.slide_layouts[0]  # 标题幻灯片
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "人工智能基础知识"
    subtitle.text = "AI入门教程"
    
    # 添加第二张幻灯片
    slide_layout = prs.slide_layouts[1]  # 标题和内容
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "什么是人工智能？"
    content.text = """人工智能（Artificial Intelligence，AI）是计算机科学的一个分支。

主要特征：
• 学习能力
• 推理能力  
• 感知能力
• 决策能力

应用领域：
• 自然语言处理
• 计算机视觉
• 机器学习
• 专家系统"""
    
    # 添加第三张幻灯片
    slide_layout = prs.slide_layouts[1]  # 标题和内容
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "AI发展历程"
    content.text = """重要里程碑：

1950年 - 图灵测试提出
1956年 - 达特茅斯会议，AI学科诞生
1980年代 - 专家系统兴起
1997年 - 深蓝击败国际象棋世界冠军
2016年 - AlphaGo击败围棋世界冠军
2020年代 - 大语言模型爆发"""
    
    # 保存文件
    filename = "test_ai_presentation.pptx"
    prs.save(filename)
    print(f"✅ 测试PPTX文件已创建: {filename}")
    return filename

def test_real_pptx_upload():
    """测试真实PPTX文件上传"""
    import requests
    import os
    
    print("\n=== 测试真实PPTX文件上传 ===")
    
    # 创建测试文件
    pptx_file = create_test_pptx()
    
    if not os.path.exists(pptx_file):
        print("❌ 测试文件创建失败")
        return
    
    try:
        BASE_URL = "http://localhost:5000"
        
        # 上传文件
        with open(pptx_file, 'rb') as f:
            files = {'file': (pptx_file, f, 'application/vnd.openxmlformats-officedocument.presentationml.presentation')}
            data = {'num_questions': '3'}
            
            response = requests.post(f"{BASE_URL}/api/quiz/upload", files=files, data=data, timeout=60)
            
            print(f"上传状态码: {response.status_code}")
            
            if response.status_code == 200:
                result = response.json()
                if result.get('success'):
                    print("✅ PPTX文件上传和题目生成成功！")
                    questions = result.get('questions', [])
                    print(f"生成了 {len(questions)} 道题目:")
                    for i, q in enumerate(questions, 1):
                        print(f"\n题目 {i}: {q.get('question', '未知题目')}")
                        print(f"原始数据: {q}")  # 添加调试信息
                        options = q.get('options', [])
                        for j, option in enumerate(options):
                            print(f"  {chr(65+j)}. {option}")
                        correct_ans = q.get('correct_answer', 0)
                        print(f"  正确答案类型: {type(correct_ans)}, 值: {correct_ans}")
                        if isinstance(correct_ans, (int, str)):
                            try:
                                if isinstance(correct_ans, str):
                                    print(f"  正确答案: {correct_ans}")
                                else:
                                    print(f"  正确答案: {chr(65+int(correct_ans))}")
                            except (ValueError, TypeError) as e:
                                print(f"  正确答案转换错误: {e}")
                        else:
                            print(f"  正确答案: 未知格式")
                else:
                    print(f"❌ 上传失败: {result.get('message')}")
            else:
                print(f"❌ 请求失败，状态码: {response.status_code}")
                print(f"错误信息: {response.text}")
                
    except requests.exceptions.RequestException as e:
        print(f"❌ 网络请求错误: {e}")
    except Exception as e:
        print(f"❌ 其他错误: {e}")
    finally:
        # 清理测试文件
        if os.path.exists(pptx_file):
            os.remove(pptx_file)
            print(f"清理测试文件: {pptx_file}")

if __name__ == "__main__":
    test_real_pptx_upload()
